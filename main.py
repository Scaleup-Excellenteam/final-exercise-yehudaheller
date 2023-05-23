import json
import collections.abc
import time
from os import path
from pptx import Presentation
import openai
import asyncio
import os
from dotenv import load_dotenv
from datetime import datetime, timedelta

# Load the environment variables from the .env file
load_dotenv()

# Get the API key from the environment variables
API_KEY = os.environ.get("API_KEY")

# Counter for the number of requests made
request_counter = 0

# Last request time
last_request_time = None

# Number of requests allowed per time interval
requests_per_interval = 3

# Time interval in minutes
interval_minutes = 1


def save_to_json(generate_text, file_name):
    """
    Saves the generated text to a JSON file.

    @param generate_text: The generated text to be saved.
    @param file_name: The name of the JSON file to save.
    """
    with open(file_name, "w") as f:
        json.dump(generate_text, f)


def read_pptx_file(file_path):
    """
    Reads the PowerPoint file and extracts the text from each slide.

    @param file_path: The path to the PowerPoint file.
    @return: A list of strings, where each string represents the text from a slide.
    """
    # Open the PowerPoint file
    prs = Presentation(file_path)

    # Initialize an empty list to hold the text from each slide
    slides_text = []

    # Loop through each slide in the presentation
    for slide in prs.slides:
        # Initialize an empty string to hold the text from the slide
        text = ""

        # Loop through each shape in the slide
        for shape in slide.shapes:
            # Check if the shape has text
            if shape.has_text_frame:
                # Loop through each paragraph in the text frame
                for paragraph in shape.text_frame.paragraphs:
                    # Add the text from the paragraph to the string
                    text += paragraph.text

        # Append the slide's text to the slides_text list
        slides_text.append(text)

    return slides_text


async def integrate_openai(slides_text):
    global request_counter
    global last_request_time

    openai.api_key = API_KEY
    model_engine = "gpt-3.5-turbo"
    generate_text = ""

    for i, slide_text in enumerate(slides_text):
        prompt = "Write a summary of the following slide:\n"
        prompt += slide_text + "\n"

        response = await asyncio.to_thread(openai.ChatCompletion.create,
                                           model=model_engine,
                                           messages=[
                                               {"role": "system", "content": prompt},
                                               {"role": "user", "content": "summarize the slide"},
                                           ]
                                           )

        request_counter += 1
        last_request_time = datetime.now()

        summary = response.choices[0].message.content.strip()
        summary = summary.replace(". ", ".\n")
        generate_text += summary + "\n"

        # print(f"test:----- current generate_text contain: ------- \n {generate_text} \n")

        if i < len(slides_text) - 1 and request_counter % requests_per_interval == 0:
            time_diff = datetime.now() - last_request_time
            remaining_time = timedelta(minutes=interval_minutes) - time_diff

            if remaining_time.total_seconds() > 0:
                await asyncio.sleep(remaining_time.total_seconds())

    return generate_text


def get_pptx_file_path():
    """
    Prompts the user to enter the path to the PowerPoint file.

    @return: The path to the PowerPoint file.
    """
    # Prompt the user to enter the path to the PowerPoint file
    filename = input("Enter the path to the PowerPoint file: ")

    # Check if the file exists
    if not path.isfile(filename):
        print("File not found. Please try again.")
        return get_pptx_file_path()  # Call the function recursively to try again

    return filename


async def main():
    """
    Main function to orchestrate the PowerPoint file processing and OpenAI integration.
    """
    # Get the path to the PowerPoint file
    pptx_file_path = get_pptx_file_path()

    # Call the read_pptx_file function to read the PowerPoint file
    slides_text = read_pptx_file(pptx_file_path)
    print(slides_text)

    # Integrate OpenAI asynchronously
    generate_text = await integrate_openai(slides_text)
    print(generate_text)

    # Get the output file name based on the PowerPoint file name
    file_name_without_extension = path.splitext(pptx_file_path)[0]
    output_file_name = file_name_without_extension + ".json"

    # Save the generated text to a JSON file with the same name as the original presentation
    save_to_json(generate_text, output_file_name)


if __name__ == "__main__":
    asyncio.run(main())
