from pptx import Presentation
from os import path


def read_pptx_file(file_path: str) -> list[str]:
    """
    Reads the PowerPoint file and extracts the text from each slide.

    Args:
        file_path (str): The path to the PowerPoint file.

    Returns:
        List[str]: A list of strings, where each string represents the text from a slide.
    """
    # Open the PowerPoint file
    prs = Presentation(file_path)

    # Initialize an empty list to hold the text from each slide, each element in list is one slide
    slides_text = []

    # Loop through each slide in the presentation
    for slide in prs.slides:
        # Initialize an empty string to hold the text from the slide
        text = ""

        # Loop through each shape in the slide
        for shape in slide.shapes:
            # Check if the shape has text
            if shape.has_text_frame:
                # Loop through each paragraph in the text frame
                for paragraph in shape.text_frame.paragraphs:
                    # Add the text from the paragraph to the string
                    text += paragraph.text

        # Append the slide's text to the slides_text list
        slides_text.append(text)

    return slides_text


def get_pptx_file_path() -> str:
    """
    Prompts the user to enter the path to the PowerPoint file.

    Returns:
        str: The path to the PowerPoint file.
    """
    # Prompt the user to enter the path to the PowerPoint file
    file_path = input("Enter the path to the PowerPoint file: ")

    # Check if the file exists
    if not path.isfile(file_path):
        print("File not found. Please try again.")
        return get_pptx_file_path()  # Call the function recursively to try again

    return file_path
